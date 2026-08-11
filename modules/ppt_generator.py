"""
Module to generate improvement action plan PPT
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from datetime import datetime

# ABB Colors
ABB_RED = RGBColor(255, 0, 15)
ABB_LILAC = RGBColor(103, 100, 246)
ABB_GRAY = RGBColor(100, 100, 100)

def generate_ppt(ppt_data, analysis, suggestions):
    """
    Generate a new PPT with improvement action plan
    
    Args:
        ppt_data: Parsed PPT data
        analysis: Analysis results
        suggestions: Generated suggestions
    
    Returns:
        str: Path to the generated PPT file
    """
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = ABB_LILAC
    
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    
    p = title_frame.paragraphs[0]
    p.text = "2026 Engagement Survey"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p = title_frame.add_paragraph()
    p.text = "Development Actions and Follow-up Plan"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = f"{analysis['department']}"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Current Status Overview
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide2.shapes.title
    title.text = "Current Engagement Status"
    
    content_box = slide2.placeholders[1].text_frame
    content_box.clear()
    
    p = content_box.paragraphs[0]
    p.text = f"Department: {analysis['department']}"
    p.font.size = Pt(16)
    p.level = 0
    
    p = content_box.add_paragraph()
    p.text = f"Overall Engagement Score: {analysis['overall_score']}%"
    p.font.size = Pt(16)
    p.level = 0
    
    p = content_box.add_paragraph()
    p.text = f"Total Participants: {analysis['total_participants']}"
    p.font.size = Pt(16)
    p.level = 0
    
    p = content_box.add_paragraph()
    p.text = f"Score Range: {analysis['score_range']['low_end']}% - {analysis['score_range']['high_end']}%"
    p.font.size = Pt(16)
    p.level = 0
    
    p = content_box.add_paragraph()
    p.text = f"Improvement Potential: {analysis['improvement_potential']}"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ABB_RED
    p.level = 0
    
    # Slide 3: Key Issues
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide3.shapes.title
    title.text = "Key Issues Identified"
    
    content_box = slide3.placeholders[1].text_frame
    content_box.clear()
    
    for i, issue in enumerate(analysis['key_issues'], 1):
        p = content_box.paragraphs[0] if i == 1 else content_box.add_paragraph()
        p.text = f"{i}. {issue}"
        p.font.size = Pt(16)
        p.level = 0
    
    p = content_box.add_paragraph()
    p.text = ""
    
    p = content_box.add_paragraph()
    p.text = "Priority Improvement Areas:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    for area in analysis['priority_areas']:
        p = content_box.add_paragraph()
        p.text = f"• {area}"
        p.font.size = Pt(14)
        p.level = 1
    
    # Slide 4: Sense of Belonging Actions
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide4.shapes.title
    title.text = "Action Plan - Sense of Belonging"
    
    content_box = slide4.placeholders[1].text_frame
    content_box.clear()
    
    p = content_box.paragraphs[0]
    p.text = f"Core Question: \"I feel a sense of belonging at ABB\""
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = ABB_RED
    p.level = 0
    
    p = content_box.add_paragraph()
    p.text = ""
    
    for i, suggestion in enumerate(suggestions['belonging'], 1):
        p = content_box.add_paragraph()
        p.text = f"{i}. {suggestion}"
        p.font.size = Pt(11)
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    # Slide 5: Work-Life Balance Actions
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide5.shapes.title
    title.text = "Action Plan - Work-Life Balance"
    
    content_box = slide5.placeholders[1].text_frame
    content_box.clear()
    
    p = content_box.paragraphs[0]
    p.text = f"Core Question: \"I am able to successfully balance my work and personal life\""
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = ABB_RED
    p.level = 0
    
    p = content_box.add_paragraph()
    p.text = ""
    
    for i, suggestion in enumerate(suggestions['work_life_balance'], 1):
        p = content_box.add_paragraph()
        p.text = f"{i}. {suggestion}"
        p.font.size = Pt(11)
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    # Slide 6: Leadership Development
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide6.shapes.title
    title.text = "Supporting Theme - Leadership Development"
    
    content_box = slide6.placeholders[1].text_frame
    content_box.clear()
    
    p = content_box.paragraphs[0]
    p.text = "Purpose: Strengthen management capabilities to effectively implement improvement initiatives"
    p.font.size = Pt(12)
    p.level = 0
    
    p = content_box.add_paragraph()
    p.text = ""
    
    p = content_box.add_paragraph()
    p.text = "Focus Areas:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    for suggestion in suggestions['leadership']:
        p = content_box.add_paragraph()
        p.text = f"• {suggestion}"
        p.font.size = Pt(12)
        p.level = 1
    
    # Slide 7: Implementation Summary
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide7.shapes.title
    title.text = "Implementation Plan Summary"
    
    content_box = slide7.placeholders[1].text_frame
    content_box.clear()
    
    p = content_box.paragraphs[0]
    p.text = f"Total Action Items: {suggestions['total_actions']} concrete improvement actions"
    p.font.size = Pt(14)
    p.font.bold = True
    p.level = 0
    
    p = content_box.add_paragraph()
    p.text = f"Sense of Belonging: {len(suggestions['belonging'])} actions"
    p.font.size = Pt(12)
    p.level = 0
    
    p = content_box.add_paragraph()
    p.text = f"Work-Life Balance: {len(suggestions['work_life_balance'])} actions"
    p.font.size = Pt(12)
    p.level = 0
    
    p = content_box.add_paragraph()
    p.text = ""
    
    p = content_box.add_paragraph()
    p.text = "Implementation Timeline:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.level = 0
    
    p = content_box.add_paragraph()
    p.text = "January 2026: Begin implementation of all action items"
    p.font.size = Pt(11)
    p.level = 1
    
    p = content_box.add_paragraph()
    p.text = "January - June 2026: Ongoing execution and monitoring"
    p.font.size = Pt(11)
    p.level = 1
    
    p = content_box.add_paragraph()
    p.text = "July 31, 2026: Completion and evaluation of all actions"
    p.font.size = Pt(11)
    p.font.bold = True
    p.level = 1
    
    # Slide 8: Expected Outcomes
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide8.shapes.title
    title.text = "Expected Outcomes & Success Metrics"
    
    content_box = slide8.placeholders[1].text_frame
    content_box.clear()
    
    p = content_box.paragraphs[0]
    p.text = "Employee Engagement Score: ↑ Increase"
    p.font.size = Pt(13)
    p.level = 0
    
    p = content_box.add_paragraph()
    p.text = "Sense of Belonging Score: ↑ Significant Increase"
    p.font.size = Pt(13)
    p.level = 0
    
    p = content_box.add_paragraph()
    p.text = "Work-Life Balance Score: ↑ Significant Increase"
    p.font.size = Pt(13)
    p.level = 0
    
    p = content_box.add_paragraph()
    p.text = "Employee Retention Rate: ↑ Improvement"
    p.font.size = Pt(13)
    p.level = 0
    
    p = content_box.add_paragraph()
    p.text = "Team Productivity: ↑ Enhancement"
    p.font.size = Pt(13)
    p.level = 0
    
    p = content_box.add_paragraph()
    p.text = "Workplace Satisfaction: ↑ Improvement"
    p.font.size = Pt(13)
    p.level = 0
    
    p = content_box.add_paragraph()
    p.text = ""
    
    p = content_box.add_paragraph()
    p.text = suggestions.get('additional_notes', '')
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = ABB_GRAY
    p.level = 0
    
    # Slide 9: Completion Date
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide9.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = ABB_RED
    
    date_box = slide9.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1.5))
    date_frame = date_box.text_frame
    date_frame.word_wrap = True
    
    p = date_frame.paragraphs[0]
    p.text = "Target Completion Date"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p = date_frame.add_paragraph()
    p.text = "July 31, 2026"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Save the presentation
    output_path = '/tmp/Engagement_Improvement_Plan.pptx'
    prs.save(output_path)
    
    return output_path
