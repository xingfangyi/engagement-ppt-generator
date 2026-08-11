"""
Module to parse employee engagement survey PPT and extract data
"""
from pptx import Presentation
import re

def parse_ppt(uploaded_file):
    """
    Parse the uploaded engagement survey PPT and extract key data
    
    Args:
        uploaded_file: The uploaded PPT file from Streamlit
    
    Returns:
        dict: Extracted data including department, scores, etc.
    """
    
    prs = Presentation(uploaded_file)
    
    data = {
        'department': 'Unknown Department',
        'overall_score': 0,
        'total_participants': 0,
        'all_scores': [],
        'lowest_scores': [],
        'lowest_items': [],
        'highest_scores': [],
        'highest_items': [],
        'survey_date': 'October 2025',
        'raw_text': [],  # Store all text for better analysis
    }
    
    # Extract all text from slides
    for slide_idx, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                text = shape.text.strip()
                slide_text.append(text)
                data['raw_text'].append(text)
        
        # Try to identify department from slide text
        slide_combined = ' '.join(slide_text)
        
        # Extract department name
        if 'Marine' in slide_combined or 'Service' in slide_combined:
            for text in slide_text:
                if 'Marine' in text or 'Service' in text:
                    data['department'] = text
                    break
        
        # Extract all percentage values
        for text in slide_text:
            matches = re.findall(r'(\d+)\s*%', text)
            if matches:
                for match in matches:
                    score = int(match)
                    if 0 <= score <= 100:
                        data['all_scores'].append(score)
        
        # Extract participant count
        participant_match = re.search(r'(\d+)\s*(?:respondent|participant|people|person)', slide_combined, re.IGNORECASE)
        if participant_match:
            data['total_participants'] = int(participant_match.group(1))
    
    # Calculate statistics
    if data['all_scores']:
        data['overall_score'] = sum(data['all_scores']) / len(data['all_scores'])
        
        # Get 10 lowest scores
        sorted_scores = sorted(data['all_scores'])
        data['lowest_scores'] = sorted_scores[:min(10, len(sorted_scores))]
        
        # Get 3 highest scores
        data['highest_scores'] = sorted_scores[-3:][::-1]
    
    return data
